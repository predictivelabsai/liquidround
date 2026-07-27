#!/usr/bin/env python3
"""Generate user guide PDF and PPTX from docs/user_guide.md.

Usage:
    python -m scripts.gen_user_guide          # generates both PDF + PPTX
    python -m scripts.gen_user_guide --pdf    # PDF only
    python -m scripts.gen_user_guide --pptx   # PPTX only
"""
from __future__ import annotations

import re
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
DOCS = ROOT / "docs"
MD_PATH = DOCS / "user_guide.md"
PDF_PATH = DOCS / "liquidround_user_guide.pdf"
PPTX_PATH = DOCS / "liquidround_user_guide.pptx"

# Brand palette
NAVY = (11, 18, 32)       # #0B1220
DARK = (17, 24, 39)       # #111827
AMBER = (245, 158, 11)    # #F59E0B
WHITE = (248, 250, 252)   # #F8FAFC
MUTED = (148, 163, 184)   # #94A3B8
SLATE = (30, 41, 59)      # #1E293B

IMG_RE = re.compile(r"^!\[([^\]]*)\]\(([^)]+)\)\s*$")


def parse_md(path: Path) -> list[dict]:
    """Parse markdown into sections: {level, title, body}."""
    text = path.read_text()
    sections = []
    current = None
    for line in text.split("\n"):
        m = re.match(r"^(#{1,3})\s+(.*)", line)
        if m:
            if current:
                sections.append(current)
            current = {"level": len(m.group(1)), "title": m.group(2).strip(), "body": ""}
        elif current:
            current["body"] += line + "\n"
    if current:
        sections.append(current)
    for s in sections:
        s["body"] = s["body"].strip()
    return sections


def _resolve_img(src: str) -> Path | None:
    """Resolve an image path relative to docs/."""
    p = DOCS / src
    return p if p.exists() else None


# ── PDF generation (reportlab) ──────────────────────────────────────

def generate_pdf(sections: list[dict], out: Path):
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.colors import Color
    from reportlab.lib.units import mm
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
        PageBreak, Image,
    )

    page_w = A4[0] - 40 * mm
    navy = Color(*[c / 255 for c in NAVY])
    amber = Color(*[c / 255 for c in AMBER])
    white_c = Color(*[c / 255 for c in WHITE])
    muted_c = Color(*[c / 255 for c in MUTED])
    dark_c = Color(*[c / 255 for c in DARK])
    slate_c = Color(*[c / 255 for c in SLATE])

    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle("LR_Title", parent=styles["Title"],
                              fontSize=28, textColor=amber, spaceAfter=6,
                              fontName="Helvetica-Bold"))
    styles.add(ParagraphStyle("LR_H1", parent=styles["Heading1"],
                              fontSize=20, textColor=amber, spaceBefore=18,
                              spaceAfter=8, fontName="Helvetica-Bold"))
    styles.add(ParagraphStyle("LR_H2", parent=styles["Heading2"],
                              fontSize=15, textColor=white_c, spaceBefore=14,
                              spaceAfter=6, fontName="Helvetica-Bold"))
    styles.add(ParagraphStyle("LR_H3", parent=styles["Heading3"],
                              fontSize=12, textColor=amber, spaceBefore=10,
                              spaceAfter=4, fontName="Helvetica-Bold"))
    styles.add(ParagraphStyle("LR_Body", parent=styles["BodyText"],
                              fontSize=10, textColor=white_c, leading=14,
                              spaceAfter=6, fontName="Helvetica"))
    styles.add(ParagraphStyle("LR_Bullet", parent=styles["BodyText"],
                              fontSize=10, textColor=white_c, leading=14,
                              leftIndent=20, bulletIndent=8, spaceAfter=3,
                              fontName="Helvetica"))
    styles.add(ParagraphStyle("LR_Caption", parent=styles["BodyText"],
                              fontSize=8, textColor=muted_c, leading=10,
                              spaceAfter=8, fontName="Helvetica-Oblique",
                              alignment=1))
    styles.add(ParagraphStyle("LR_TableCell", fontSize=9, textColor=white_c,
                              fontName="Helvetica", leading=11))
    styles.add(ParagraphStyle("LR_TableHeader", fontSize=9, textColor=navy,
                              fontName="Helvetica-Bold", leading=11))

    def on_page(canvas, doc):
        canvas.saveState()
        canvas.setFillColor(navy)
        canvas.rect(0, 0, A4[0], A4[1], fill=1, stroke=0)
        canvas.setFillColor(muted_c)
        canvas.setFont("Helvetica", 7)
        canvas.drawString(20 * mm, 10 * mm, "LiquidRound v0.5.0 — Predictive Labs Ltd")
        canvas.drawRightString(A4[0] - 20 * mm, 10 * mm, f"Page {doc.page}")
        canvas.restoreState()

    doc = SimpleDocTemplate(str(out), pagesize=A4,
                            leftMargin=20 * mm, rightMargin=20 * mm,
                            topMargin=20 * mm, bottomMargin=20 * mm)
    story = []

    # Cover page
    story.append(Spacer(1, 60 * mm))
    story.append(Paragraph("LiquidRound", styles["LR_Title"]))
    story.append(Paragraph("User Guide v0.5.0", styles["LR_H2"]))
    story.append(Spacer(1, 10 * mm))
    story.append(Paragraph(
        "AI-powered M&A, IPO readiness, and hedge fund intelligence platform",
        styles["LR_Body"]))
    story.append(Paragraph("Predictive Labs Ltd", styles["LR_Body"]))
    story.append(PageBreak())

    def _add_image(src: str, caption: str):
        img_path = _resolve_img(src)
        if not img_path:
            return
        try:
            img = Image(str(img_path), width=page_w, height=page_w * 0.56)
            img.hAlign = "CENTER"
            story.append(Spacer(1, 3 * mm))
            story.append(img)
            if caption:
                story.append(Paragraph(caption, styles["LR_Caption"]))
            story.append(Spacer(1, 2 * mm))
        except Exception as e:
            print(f"  Warning: could not embed {src}: {e}")

    def _render_body(body: str):
        # Parse markdown table
        table_lines = [l for l in body.split("\n") if l.startswith("|")]
        if len(table_lines) >= 3:
            rows = []
            for i, tl in enumerate(table_lines):
                if set(tl.replace("|", "").replace("-", "").strip()) <= {" ", ""}:
                    continue
                cells = [c.strip().replace("`", "") for c in tl.split("|")[1:-1]]
                if i == 0:
                    rows.append([Paragraph(c, styles["LR_TableHeader"]) for c in cells])
                else:
                    rows.append([Paragraph(c, styles["LR_TableCell"]) for c in cells])
            if rows:
                ncols = len(rows[0])
                col_w = page_w / ncols
                t = Table(rows, colWidths=[col_w] * ncols)
                t.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), amber),
                    ("BACKGROUND", (0, 1), (-1, -1), dark_c),
                    ("TEXTCOLOR", (0, 0), (-1, 0), navy),
                    ("TEXTCOLOR", (0, 1), (-1, -1), white_c),
                    ("GRID", (0, 0), (-1, -1), 0.5, slate_c),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("TOPPADDING", (0, 0), (-1, -1), 4),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                    ("LEFTPADDING", (0, 0), (-1, -1), 6),
                ]))
                story.append(t)
                story.append(Spacer(1, 4 * mm))
            non_table = [l for l in body.split("\n") if not l.startswith("|")]
            body = "\n".join(non_table).strip()
            if not body:
                return

        for line in body.split("\n"):
            line = line.strip()
            if not line:
                story.append(Spacer(1, 2 * mm))
                continue
            if line.startswith("---"):
                continue
            if line.startswith("> "):
                text = line[2:]
                text = re.sub(r"\*\*(.*?)\*\*", r"<b>\1</b>", text)
                story.append(Paragraph(f"<i>{text}</i>", styles["LR_Body"]))
                continue
            img_m = IMG_RE.match(line)
            if img_m:
                _add_image(img_m.group(2), img_m.group(1))
                continue
            if line.startswith("- **"):
                text = line[2:]
                text = re.sub(r"\*\*(.*?)\*\*", r"<b>\1</b>", text)
                text = re.sub(r"`(.*?)`", r'<font face="Courier" color="#F59E0B">\1</font>', text)
                story.append(Paragraph(f"• {text}", styles["LR_Bullet"]))
            elif line.startswith("- "):
                text = line[2:]
                text = re.sub(r"\*\*(.*?)\*\*", r"<b>\1</b>", text)
                text = re.sub(r"`(.*?)`", r'<font face="Courier" color="#F59E0B">\1</font>', text)
                story.append(Paragraph(f"• {text}", styles["LR_Bullet"]))
            else:
                text = re.sub(r"\*\*(.*?)\*\*", r"<b>\1</b>", line)
                text = re.sub(r"`(.*?)`", r'<font face="Courier" color="#F59E0B">\1</font>', text)
                story.append(Paragraph(text, styles["LR_Body"]))

    for sec in sections:
        if sec["level"] == 1 and sec["title"] != "LiquidRound":
            story.append(Paragraph(sec["title"], styles["LR_H1"]))
        elif sec["level"] == 1:
            continue
        elif sec["level"] == 2:
            story.append(Paragraph(sec["title"], styles["LR_H2"]))
        elif sec["level"] == 3:
            story.append(Paragraph(sec["title"], styles["LR_H3"]))

        if sec["body"]:
            _render_body(sec["body"])

    doc.build(story, onFirstPage=on_page, onLaterPages=on_page)
    print(f"PDF: {out} ({out.stat().st_size // 1024} KB)")


# ── PPTX generation (python-pptx) ──────────────────────────────────

def generate_pptx(sections: list[dict], out: Path):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    navy_rgb = RGBColor(*NAVY)
    dark_rgb = RGBColor(*DARK)
    amber_rgb = RGBColor(*AMBER)
    white_rgb = RGBColor(*WHITE)
    muted_rgb = RGBColor(*MUTED)

    def set_slide_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_text_box(slide, left, top, width, height, text, font_size=18,
                     bold=False, color=white_rgb, alignment=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = alignment
        return tf

    def add_image(slide, src: str, left, top, width):
        img_path = _resolve_img(src)
        if not img_path:
            return top
        try:
            slide.shapes.add_picture(str(img_path), left, top, width=width)
            from PIL import Image as PILImage
            with PILImage.open(img_path) as im:
                aspect = im.height / im.width
            return top + int(width * aspect) + Inches(0.2)
        except Exception as e:
            print(f"  Warning: could not embed {src}: {e}")
            return top

    # Cover slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, navy_rgb)
    add_text_box(slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
                 "LiquidRound", font_size=48, bold=True, color=amber_rgb,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
                 "User Guide v0.5.0", font_size=28, color=white_rgb,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
                 "AI-powered M&A, IPO readiness, and hedge fund intelligence",
                 font_size=16, color=muted_rgb, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
                 "Predictive Labs Ltd", font_size=14, color=muted_rgb,
                 alignment=PP_ALIGN.CENTER)

    # Group sections into slides by H2
    slide_groups = []
    for sec in sections:
        if sec["level"] == 1:
            continue
        if sec["level"] == 2:
            slide_groups.append({"h2": sec, "h3s": []})
        elif sec["level"] == 3 and slide_groups:
            slide_groups[-1]["h3s"].append(sec)

    for group in slide_groups:
        h2 = group["h2"]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide, navy_rgb)

        # Section title
        add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.7),
                     h2["title"], font_size=28, bold=True, color=amber_rgb)

        # Collect all content lines and images
        all_content = []
        if h2["body"]:
            all_content.append(("body", h2["body"]))
        for h3 in group["h3s"]:
            all_content.append(("h3", h3["title"]))
            if h3["body"]:
                all_content.append(("body", h3["body"]))

        # Find images in content
        images = []
        text_parts = []
        for ctype, content in all_content:
            if ctype == "body":
                for line in content.split("\n"):
                    img_m = IMG_RE.match(line.strip())
                    if img_m:
                        images.append(img_m.group(2))
                    else:
                        text_parts.append((ctype, line))
            else:
                text_parts.append((ctype, content))

        combined_body = h2["body"] + "\n" + "\n".join(
            h3["body"] for h3 in group["h3s"])
        table_lines = [l for l in combined_body.split("\n") if l.startswith("|")]

        if images:
            # Layout: text on left, image on right (or image full-width if little text)
            has_text = any(
                l.strip() and not l.strip().startswith("---") and not IMG_RE.match(l.strip())
                for _, l in text_parts if _ == "body"
            ) or any(ct == "h3" for ct, _ in text_parts)

            if has_text:
                # Two-column: text left, image right
                img_left = Inches(6.8)
                img_top = Inches(1.1)
                img_w = Inches(6.0)
                txt_w = Inches(5.8)
            else:
                # Full-width image
                img_left = Inches(1.5)
                img_top = Inches(1.1)
                img_w = Inches(10.0)
                txt_w = Inches(11.5)

            for img_src in images[:1]:
                add_image(slide, img_src, img_left, img_top, img_w)

            # Render text
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.1),
                                             txt_w, Inches(5.8))
            tf = txBox.text_frame
            tf.word_wrap = True
            first_para = True
            for ctype, content in text_parts:
                if ctype == "h3":
                    p = tf.paragraphs[0] if first_para else tf.add_paragraph()
                    first_para = False
                    p.text = content
                    p.font.size = Pt(16)
                    p.font.bold = True
                    p.font.color.rgb = amber_rgb
                    p.space_before = Pt(10)
                    p.space_after = Pt(4)
                elif ctype == "body":
                    line = content.strip()
                    if not line or line.startswith("---"):
                        continue
                    if IMG_RE.match(line):
                        continue
                    p = tf.paragraphs[0] if first_para else tf.add_paragraph()
                    first_para = False
                    clean = re.sub(r"\*\*(.*?)\*\*", r"\1", line)
                    clean = re.sub(r"`(.*?)`", r"\1", clean)
                    if line.startswith("- "):
                        clean = "• " + clean[2:]
                        p.level = 1
                    if line.startswith("> "):
                        clean = clean[2:]
                    p.text = clean
                    p.font.size = Pt(13)
                    p.font.color.rgb = white_rgb
                    p.space_after = Pt(3)

        elif len(table_lines) >= 3:
            rows_data = []
            for tl in table_lines:
                if set(tl.replace("|", "").replace("-", "").strip()) <= {" ", ""}:
                    continue
                cells = [c.strip().replace("`", "") for c in tl.split("|")[1:-1]]
                rows_data.append(cells)

            if rows_data:
                ncols = len(rows_data[0])
                nrows = len(rows_data)
                tbl_w = Inches(11.5)
                tbl_h = Inches(min(nrows * 0.35, 5.5))
                tbl = slide.shapes.add_table(nrows, ncols,
                                             Inches(0.8), Inches(1.1), tbl_w, tbl_h).table
                col_w = tbl_w // ncols
                for ci in range(ncols):
                    tbl.columns[ci].width = col_w
                for ri, row in enumerate(rows_data):
                    for ci, cell in enumerate(row):
                        c = tbl.cell(ri, ci)
                        c.text = cell
                        p = c.text_frame.paragraphs[0]
                        p.font.size = Pt(10 if ri > 0 else 11)
                        p.font.bold = ri == 0
                        p.font.color.rgb = navy_rgb if ri == 0 else white_rgb
                        c.fill.solid()
                        c.fill.fore_color.rgb = amber_rgb if ri == 0 else dark_rgb
        else:
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.1),
                                             Inches(11.5), Inches(5.8))
            tf = txBox.text_frame
            tf.word_wrap = True
            first_para = True
            for ctype, content in all_content:
                if ctype == "h3":
                    p = tf.paragraphs[0] if first_para else tf.add_paragraph()
                    first_para = False
                    p.text = content
                    p.font.size = Pt(18)
                    p.font.bold = True
                    p.font.color.rgb = amber_rgb
                    p.space_before = Pt(12)
                    p.space_after = Pt(4)
                elif ctype == "body":
                    for line in content.split("\n"):
                        line = line.strip()
                        if not line or line.startswith("---"):
                            continue
                        if IMG_RE.match(line):
                            continue
                        p = tf.paragraphs[0] if first_para else tf.add_paragraph()
                        first_para = False
                        clean = re.sub(r"\*\*(.*?)\*\*", r"\1", line)
                        clean = re.sub(r"`(.*?)`", r"\1", clean)
                        if line.startswith("- "):
                            clean = "• " + clean[2:]
                            p.level = 1
                        if line.startswith("> "):
                            clean = clean[2:]
                        p.text = clean
                        p.font.size = Pt(14)
                        p.font.color.rgb = white_rgb
                        p.space_after = Pt(4)

    prs.save(str(out))
    print(f"PPTX: {out} ({out.stat().st_size // 1024} KB)")


if __name__ == "__main__":
    import argparse
    ap = argparse.ArgumentParser(description="Generate user guide PDF/PPTX")
    ap.add_argument("--pdf", action="store_true", help="PDF only")
    ap.add_argument("--pptx", action="store_true", help="PPTX only")
    args = ap.parse_args()

    sections = parse_md(MD_PATH)
    do_both = not args.pdf and not args.pptx

    if do_both or args.pdf:
        generate_pdf(sections, PDF_PATH)
    if do_both or args.pptx:
        generate_pptx(sections, PPTX_PATH)
