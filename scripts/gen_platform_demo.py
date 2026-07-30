#!/usr/bin/env python3
"""Archive prior guide outputs and build a timestamped LiquidRound demo PDF + PPTX.

The PDF is rendered via pandoc → HTML → WeasyPrint using docs/assets/guide.css
for a dark navy slide-deck look on A4 landscape.  The PPTX is a 16:9 branded
slide deck with one slide per H2 section, screenshots floated on the right,
and native tables.
"""
from __future__ import annotations

import argparse
import logging
import re
import shutil
import subprocess
from datetime import datetime, timezone
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
DOCS = ROOT / "docs"
SOURCE = DOCS / "user_guide.md"
CSS = DOCS / "assets" / "guide.css"
ARCHIVE = DOCS / "archive"
PREFIX = "liquidround-ai-platform-demo"

log = logging.getLogger(__name__)

# ── Brand palette (RGB tuples for python-pptx) ──────────────────────────
_NAVY = (11, 18, 32)
_DARK = (17, 24, 39)
_AMBER = (245, 158, 11)
_WHITE = (248, 250, 252)
_MUTED = (148, 163, 184)

IMG_RE = re.compile(r"^!\[([^\]]*)\]\(([^)]+)\)\s*$")


def timestamp_slug() -> str:
    return datetime.now(timezone.utc).strftime("%Y%m%dT%H%M%SZ")


def archive_outputs(stamp: str) -> list[Path]:
    """Move prior timestamped demo outputs into one timestamped archive folder."""
    targets = sorted(
        p for p in DOCS.glob(f"{PREFIX}-*")
        if p.is_file() and p.suffix.lower() in {".md", ".pdf", ".docx", ".pptx", ".html"}
    )
    if not targets:
        return []
    destination = ARCHIVE / stamp
    destination.mkdir(parents=True, exist_ok=True)
    moved = []
    for path in targets:
        dest = destination / path.name
        collision = 1
        while dest.exists():
            dest = destination / f"{path.stem}-archived-{collision}{path.suffix}"
            collision += 1
        shutil.move(str(path), str(dest))
        moved.append(dest)
    return moved


def parse_source(path: Path) -> tuple[str, list[dict]]:
    text = path.read_text(encoding="utf-8")
    title = "LiquidRound AI Platform Demo"
    sections: list[dict] = []
    current = None
    for line in text.splitlines():
        match = re.match(r"^(#{1,3})\s+(.+)$", line)
        if match:
            level = len(match.group(1))
            heading = match.group(2).strip()
            if level == 1:
                title = heading
                continue
            if current:
                sections.append(current)
            current = {"level": level, "title": heading, "lines": []}
        elif current is not None:
            current["lines"].append(line)
    if current:
        sections.append(current)
    return title, sections


# ── PDF via pandoc + WeasyPrint ──────────────────────────────────────────

def build_pdf(source: Path, output: Path, stamp: str) -> None:
    """Render markdown → HTML (pandoc) → PDF (WeasyPrint) with dark CSS."""

    # The markdown already contains a cover-page div — strip the H1 title so
    # pandoc does not duplicate it. Render once to discover stable section
    # pages, then render the final document with a page-numbered TOC.
    raw = source.read_text(encoding="utf-8")
    raw = re.sub(r"^# .+\n+", "", raw)

    tmp_md = DOCS / f".tmp-{stamp}.md"
    tmp_html = DOCS / f".tmp-{stamp}.html"

    def render(markdown: str) -> None:
        tmp_md.write_text(markdown, encoding="utf-8")
        css_rel = CSS.relative_to(DOCS)
        pandoc_cmd = [
            "pandoc", str(tmp_md), "--standalone", f"--css={css_rel}",
            "--metadata", "title=LiquidRound AI Platform Demo",
            "--from=markdown+raw_html", "--to=html5", "-o", str(tmp_html),
        ]
        result = subprocess.run(
            pandoc_cmd, capture_output=True, text=True, cwd=ROOT
        )
        if result.returncode != 0:
            raise RuntimeError(f"pandoc failed: {result.stderr}")
        result = subprocess.run(
            ["weasyprint", str(tmp_html), str(output), "--base-url", str(DOCS)],
            capture_output=True, text=True,
        )
        if result.returncode != 0:
            raise RuntimeError(f"WeasyPrint failed: {result.stderr}")

    try:
        render(raw)
        info = subprocess.run(
            ["pdfinfo", str(output)], capture_output=True, text=True, check=True
        )
        page_count_match = re.search(r"^Pages:\s+(\d+)$", info.stdout, re.MULTILINE)
        page_count = int(page_count_match.group(1)) if page_count_match else 0
        page_text = []
        for page_number in range(1, page_count + 1):
            extracted = subprocess.run(
                [
                    "pdftotext", "-f", str(page_number), "-l", str(page_number),
                    str(output), "-",
                ],
                capture_output=True, text=True, check=True,
            )
            page_text.append(extracted.stdout)
        _, sections = parse_source(source)
        headings = [section["title"] for section in sections if section["level"] == 2]
        toc_lines = []
        for heading in headings:
            draft_page = next(
                (
                    index + 1
                    for index, text in enumerate(page_text)
                    if heading in {line.strip() for line in text.splitlines()}
                ),
                None,
            )
            if draft_page is not None:
                toc_lines.append(f"- {heading} — Page {draft_page + 1}")
        toc = "## Table of Contents\n\n" + "\n".join(toc_lines)
        marker = "\n---\n"
        full_md = raw.replace(marker, f"{marker}\n{toc}\n{marker}", 1)
        render(full_md)
        log.info("Generated page-numbered PDF: %s", output)
    finally:
        tmp_md.unlink(missing_ok=True)
        tmp_html.unlink(missing_ok=True)


# ── PPTX via python-pptx ─────────────────────────────────────────────────

def build_pptx(source: Path, output: Path, stamp: str) -> None:
    """Build a 16:9 branded PPTX slide deck from the markdown source."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    title, sections = parse_source(source)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    navy = RGBColor(*_NAVY)
    dark = RGBColor(*_DARK)
    amber = RGBColor(*_AMBER)
    white = RGBColor(*_WHITE)
    muted = RGBColor(*_MUTED)

    def _bg(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _text(slide, left, top, w, h, text, size=18, bold=False,
              color=white, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(left, top, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
        return tf

    def _resolve_img(src: str) -> Path | None:
        p = (DOCS / src).resolve() if not src.startswith("/") else Path(src)
        return p if p.exists() else None

    def _add_image(slide, src, left, top, width):
        img_path = _resolve_img(src)
        if not img_path:
            return top
        try:
            slide.shapes.add_picture(str(img_path), left, top, width=width)
            from PIL import Image as PILImage
            with PILImage.open(str(img_path)) as im:
                aspect = im.height / im.width
            return top + int(width * aspect) + Inches(0.15)
        except Exception as e:
            log.warning("pptx: could not embed %s: %s", src, e)
            return top

    def _clean(line: str) -> str:
        line = re.sub(r"\*\*(.*?)\*\*", r"\1", line)
        line = re.sub(r"`(.*?)`", r"\1", line)
        return line

    # ── Cover slide ─────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, navy)
    _text(slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
          title, size=44, bold=True, color=amber, align=PP_ALIGN.CENTER)
    _text(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
          "M&A · IPO · Public Markets · Investor Relations",
          size=24, color=white, align=PP_ALIGN.CENTER)
    _text(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
          "32 specialist agents · 7 categories · chat-first workspace",
          size=16, color=muted, align=PP_ALIGN.CENTER)
    _text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
          f"Predictive Labs Ltd · {stamp}",
          size=14, color=muted, align=PP_ALIGN.CENTER)

    # ── Group sections by H2 ────────────────────────────────────────
    slide_groups: list[dict] = []
    for sec in sections:
        if sec["level"] == 1:
            continue
        if sec["level"] == 2:
            slide_groups.append({"h2": sec, "h3s": []})
        elif sec["level"] == 3 and slide_groups:
            slide_groups[-1]["h3s"].append(sec)

    for group in slide_groups:
        h2 = group["h2"]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _bg(slide, navy)

        _text(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.7),
              h2["title"], size=28, bold=True, color=amber)

        # Collect content lines and images from H2 body + H3 bodies
        all_content: list[tuple[str, str]] = []
        if h2["lines"]:
            all_content.append(("body", "\n".join(h2["lines"])))
        for h3 in group["h3s"]:
            all_content.append(("h3", h3["title"]))
            if h3["lines"]:
                all_content.append(("body", "\n".join(h3["lines"])))

        images: list[str] = []
        text_parts: list[tuple[str, str]] = []
        for ctype, content in all_content:
            if ctype == "body":
                for line in content.split("\n"):
                    img_m = IMG_RE.match(line.strip())
                    if img_m:
                        images.append(img_m.group(2))
                    else:
                        text_parts.append((ctype, line))
            else:
                text_parts.append((ctype, content))

        combined = "\n".join(c for _, c in all_content if _ == "body")
        table_lines = [l for l in combined.split("\n") if l.startswith("|")]

        if images:
            has_text = any(
                l.strip() and not l.strip().startswith("---") and not IMG_RE.match(l.strip())
                for _, l in text_parts if _ == "body"
            ) or any(ct == "h3" for ct, _ in text_parts)

            if has_text:
                img_left, img_top, img_w, txt_w = Inches(6.8), Inches(1.1), Inches(6.0), Inches(5.8)
            else:
                img_left, img_top, img_w, txt_w = Inches(1.5), Inches(1.1), Inches(10.0), Inches(11.5)

            for img_src in images[:1]:
                _add_image(slide, img_src, img_left, img_top, img_w)

            tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), txt_w, Inches(5.8))
            tf = tb.text_frame
            tf.word_wrap = True
            first = True
            for ctype, content in text_parts:
                if ctype == "h3":
                    p = tf.paragraphs[0] if first else tf.add_paragraph()
                    first = False
                    p.text = content
                    p.font.size = Pt(16)
                    p.font.bold = True
                    p.font.color.rgb = amber
                    p.space_before = Pt(10)
                    p.space_after = Pt(4)
                elif ctype == "body":
                    line = content.strip()
                    if not line or line.startswith("---"):
                        continue
                    if IMG_RE.match(line):
                        continue
                    p = tf.paragraphs[0] if first else tf.add_paragraph()
                    first = False
                    clean = _clean(line)
                    if line.startswith("- "):
                        clean = "• " + clean[2:]
                        p.level = 1
                    if line.startswith("> "):
                        clean = clean[2:]
                    p.text = clean
                    p.font.size = Pt(13)
                    p.font.color.rgb = white
                    p.space_after = Pt(3)

        elif len(table_lines) >= 3:
            rows_data = []
            for tl in table_lines:
                if set(tl.replace("|", "").replace("-", "").strip()) <= {" ", ""}:
                    continue
                cells = [c.strip().replace("`", "") for c in tl.split("|")[1:-1]]
                rows_data.append(cells)
            if rows_data:
                ncols = len(rows_data[0])
                nrows = len(rows_data)
                tbl_w = Inches(11.5)
                tbl_h = Inches(min(nrows * 0.35, 5.5))
                tbl = slide.shapes.add_table(nrows, ncols,
                                             Inches(0.8), Inches(1.1), tbl_w, tbl_h).table
                col_w = tbl_w // ncols
                for ci in range(ncols):
                    tbl.columns[ci].width = col_w
                for ri, row in enumerate(rows_data):
                    for ci, cell in enumerate(row):
                        c = tbl.cell(ri, ci)
                        c.text = cell
                        p = c.text_frame.paragraphs[0]
                        p.font.size = Pt(10 if ri > 0 else 11)
                        p.font.bold = ri == 0
                        p.font.color.rgb = navy if ri == 0 else white
                        c.fill.solid()
                        c.fill.fore_color.rgb = amber if ri == 0 else dark
        else:
            tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.5), Inches(5.8))
            tf = tb.text_frame
            tf.word_wrap = True
            first = True
            for ctype, content in all_content:
                if ctype == "h3":
                    p = tf.paragraphs[0] if first else tf.add_paragraph()
                    first = False
                    p.text = content
                    p.font.size = Pt(18)
                    p.font.bold = True
                    p.font.color.rgb = amber
                    p.space_before = Pt(12)
                    p.space_after = Pt(4)
                elif ctype == "body":
                    for line in content.split("\n"):
                        line = line.strip()
                        if not line or line.startswith("---"):
                            continue
                        if IMG_RE.match(line):
                            continue
                        p = tf.paragraphs[0] if first else tf.add_paragraph()
                        first = False
                        clean = _clean(line)
                        if line.startswith("- "):
                            clean = "• " + clean[2:]
                            p.level = 1
                        if line.startswith("> "):
                            clean = clean[2:]
                        p.text = clean
                        p.font.size = Pt(14)
                        p.font.color.rgb = white
                        p.space_after = Pt(4)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--timestamp", default=timestamp_slug())
    parser.add_argument("--no-archive", action="store_true")
    parser.add_argument("--pdf", action="store_true", help="PDF only")
    parser.add_argument("--pptx", action="store_true", help="PPTX only")
    args = parser.parse_args()
    stamp = args.timestamp
    if not re.fullmatch(r"\d{8}T\d{6}Z", stamp):
        parser.error("--timestamp must use YYYYMMDDTHHMMSSZ")
    if not SOURCE.exists():
        parser.error(f"missing source: {SOURCE}")

    moved = [] if args.no_archive else archive_outputs(stamp)
    slug = f"{PREFIX}-{stamp}"
    markdown_out = DOCS / f"{slug}.md"
    pdf_out = DOCS / f"{slug}.pdf"
    pptx_out = DOCS / f"{slug}.pptx"
    shutil.copy2(SOURCE, markdown_out)
    do_both = not args.pdf and not args.pptx
    if do_both or args.pdf:
        build_pdf(SOURCE, pdf_out, stamp)
        print(f"PDF: {pdf_out} ({pdf_out.stat().st_size // 1024} KB)")
    if do_both or args.pptx:
        build_pptx(SOURCE, pptx_out, stamp)
        print(f"PPTX: {pptx_out} ({pptx_out.stat().st_size // 1024} KB)")

    print(f"Archived: {len(moved)} file(s)")
    print(f"Markdown: {markdown_out}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
