"""
Document parsing utilities for XLS, PPT, and PDF uploads.
"""
import os
from pathlib import Path
from typing import Optional


class DocumentParser:
    """Parse uploaded documents and extract key M&A-relevant data."""

    def parse(self, file_path: str) -> dict:
        ext = Path(file_path).suffix.lower()
        if ext in (".xlsx", ".xls"):
            return self.parse_xlsx(file_path)
        elif ext in (".pptx", ".ppt"):
            return self.parse_pptx(file_path)
        elif ext == ".pdf":
            return self.parse_pdf(file_path)
        return {"error": f"Unsupported file type: {ext}"}

    def parse_xlsx(self, file_path: str) -> dict:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)
        sheets = {}
        for name in wb.sheetnames:
            ws = wb[name]
            rows = []
            for row in ws.iter_rows(max_row=200, values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append(cells)
            sheets[name] = rows
        return {
            "type": "xlsx",
            "filename": os.path.basename(file_path),
            "sheets": sheets,
            "sheet_names": wb.sheetnames,
            "summary": f"Excel workbook with {len(wb.sheetnames)} sheet(s): {', '.join(wb.sheetnames)}",
        }

    def parse_pptx(self, file_path: str) -> dict:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
            slides.append({"slide_number": i + 1, "texts": texts})
        return {
            "type": "pptx",
            "filename": os.path.basename(file_path),
            "slides": slides,
            "slide_count": len(slides),
            "summary": f"PowerPoint with {len(slides)} slide(s)",
        }

    def parse_pdf(self, file_path: str) -> dict:
        if not os.path.exists(file_path):
            return {"error": f"File not found: {file_path}"}
        import pdfplumber
        pages = []
        tables_found = []
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages[:50]):
                text = page.extract_text() or ""
                page_tables = page.extract_tables() or []
                pages.append({"page_number": i + 1, "text": text[:2000]})
                for t in page_tables:
                    tables_found.append({"page": i + 1, "rows": len(t), "data": t[:20]})
        return {
            "type": "pdf",
            "filename": os.path.basename(file_path),
            "pages": pages,
            "page_count": len(pages),
            "tables": tables_found,
            "summary": f"PDF with {len(pages)} page(s), {len(tables_found)} table(s)",
        }

    def extract_all_text(self, parsed: dict) -> str:
        """Flatten parsed document into a single text string for LLM consumption."""
        parts = []
        doc_type = parsed.get("type", "")
        if doc_type == "xlsx":
            for name, rows in parsed.get("sheets", {}).items():
                parts.append(f"=== Sheet: {name} ===")
                for row in rows[:100]:
                    parts.append(" | ".join(row))
        elif doc_type == "pptx":
            for slide in parsed.get("slides", []):
                parts.append(f"=== Slide {slide['slide_number']} ===")
                parts.extend(slide["texts"])
        elif doc_type == "pdf":
            for page in parsed.get("pages", []):
                parts.append(f"=== Page {page['page_number']} ===")
                parts.append(page["text"])
        return "\n".join(parts)


document_parser = DocumentParser()
