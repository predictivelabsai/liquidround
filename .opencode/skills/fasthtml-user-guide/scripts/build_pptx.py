#!/usr/bin/env python3
"""Build a PowerPoint (.pptx) slide deck from a FastClinic user-guide md.

Slides are separated by `---`; each leads with a heading and may embed a
screenshot (`![alt](img/..)`) and/or markdown tables. The PPTX is kept visually
**in sync with the PDF**:
  - soft-wrapped markdown lines are joined into one paragraph/bullet,
  - inline markdown (`**bold**`, `*italic*`, `` `code` ``, links) renders as
    styled runs (not raw markers),
  - markdown tables become native PowerPoint tables,
  - a leading `# H1` slide (optionally wrapped in a `::: cover` fenced div) is
    rendered as a branded cover page.

Usage: python scripts/build_pptx.py <input.md> <output.pptx> [deck title]
"""
from __future__ import annotations

import os
import re
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# FastClinic palette
TEAL = RGBColor(0x1E, 0x6F, 0xB8)      # FastClinic blue
TEAL_DK = RGBColor(0x17, 0x56, 0x8F)   # darker blue
DARK = RGBColor(0x1B, 0x27, 0x33)
MUTE = RGBColor(0x6B, 0x72, 0x78)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ROW_ALT = RGBColor(0xEE, 0xF4, 0xFA)   # blue tint
CREAM = RGBColor(0xD6, 0xE6, 0xF5)     # pale blue

FOOTER = "FastClinic · Modern primary care, made personal · fastclinic.example"

IMG_RE = re.compile(r"!\[(?P<alt>[^\]]*)\]\((?P<path>[^)]+)\)")
HEAD_RE = re.compile(r"^(#{1,6})\s+(.*)$")
BULLET_RE = re.compile(r"^\s*([-*]|\d+\.)\s+(.*)$")
SEP_ROW_RE = re.compile(r"^\s*\|?[\s:|-]+\|?\s*$")
INLINE_RE = re.compile(r"(\*\*.+?\*\*|`.+?`|\*.+?\*|\[[^\]]+\]\([^)]+\))")


def _split_row(line: str) -> list[str]:
    s = line.strip()
    if s.startswith("|"):
        s = s[1:]
    if s.endswith("|"):
        s = s[:-1]
    return [c.strip() for c in s.split("|")]


def _plain(text: str) -> str:
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    return text.strip()


def parse_slides(md: str) -> list[dict]:
    """Each slide → {title, level, image, blocks}; blocks keep raw inline markdown.
    Soft-wrapped continuation lines are joined into their paragraph/bullet."""
    slides = []
    for raw in re.split(r"(?m)^---\s*$", md):
        lines = raw.split("\n")
        title, level, image, blocks, para = "", 0, None, [], []

        def flush():
            nonlocal para
            if para:
                blocks.append(("p", " ".join(para).strip()))
                para = []

        i = 0
        while i < len(lines):
            ln = lines[i]
            s = ln.strip()
            if s.startswith(":::"):           # fenced-div markers (cover) — skip
                flush(); i += 1; continue
            if not s:                          # blank line ends a paragraph
                flush(); i += 1; continue
            if s.startswith("|"):              # table block
                flush(); rows = []
                while i < len(lines) and lines[i].strip().startswith("|"):
                    if not SEP_ROW_RE.match(lines[i]):
                        rows.append(_split_row(lines[i]))
                    i += 1
                if rows:
                    blocks.append(("table", rows))
                continue
            mimg = IMG_RE.search(s)
            if mimg:
                flush()
                if image is None:
                    image = mimg.group("path").strip()
                i += 1; continue
            mhead = HEAD_RE.match(s)
            if mhead:
                flush()
                if not title:
                    title, level = mhead.group(2).strip(), len(mhead.group(1))
                else:
                    blocks.append(("p", mhead.group(2).strip()))
                i += 1; continue
            mb = BULLET_RE.match(ln)
            if mb:
                flush(); para = ["• " + mb.group(2).strip()]
                i += 1; continue
            para.append(s.lstrip("> ").strip())  # continuation / paragraph line
            i += 1
        flush()
        if title or image or blocks:
            slides.append({"title": title, "level": level, "image": image, "blocks": blocks})
    return slides


def _add_runs(para, text, size, color):
    """Render text with inline markdown as styled runs (bold/italic/code/link)."""
    for tok in re.split(INLINE_RE, text):
        if not tok:
            continue
        run = para.add_run()
        if tok.startswith("**") and tok.endswith("**") and len(tok) > 4:
            run.text = tok[2:-2]; run.font.bold = True; run.font.color.rgb = TEAL_DK
        elif tok.startswith("`") and tok.endswith("`") and len(tok) > 2:
            run.text = tok[1:-1]; run.font.name = "Consolas"; run.font.color.rgb = TEAL_DK
        elif tok.startswith("*") and tok.endswith("*") and len(tok) > 2:
            run.text = tok[1:-1]; run.font.italic = True; run.font.color.rgb = color
        elif tok.startswith("["):
            m = re.match(r"\[([^\]]+)\]\([^)]+\)", tok)
            run.text = m.group(1) if m else tok; run.font.color.rgb = TEAL
        else:
            run.text = tok; run.font.color.rgb = color
        run.font.size = size


def _add_title(slide, text, sw):
    band = slide.shapes.add_shape(1, 0, 0, sw, Inches(1.0))
    band.fill.solid(); band.fill.fore_color.rgb = TEAL; band.line.fill.background()
    tf = band.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.45); tf.margin_top = Inches(0.14)
    r = tf.paragraphs[0].add_run(); r.text = text
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = WHITE


def _add_text(slide, blocks, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame; tf.word_wrap = True
    first = True
    for kind, val in blocks:
        if kind != "p":
            continue
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        _add_runs(para, val, Pt(14), DARK)
        para.space_after = Pt(6)


def _add_table(slide, rows, left, top, width, height):
    nrows, ncols = len(rows), max(len(r) for r in rows)
    rows = [r + [""] * (ncols - len(r)) for r in rows]
    gf = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    tbl = gf.table
    if ncols == 3:
        tbl.columns[0].width = int(width * 0.26)
        tbl.columns[1].width = int(width * 0.16)
        tbl.columns[2].width = int(width * 0.58)
    body_pt = 10 if nrows <= 16 else (8 if nrows <= 26 else 7)
    for r in range(nrows):
        for c in range(ncols):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.01); cell.margin_bottom = Inches(0.01)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = TEAL if r == 0 else (ROW_ALT if r % 2 == 0 else WHITE)
            para = cell.text_frame.paragraphs[0]
            if r == 0:
                run = para.add_run(); run.text = _plain(rows[r][c])
                run.font.size = Pt(body_pt); run.font.bold = True; run.font.color.rgb = WHITE
            else:
                _add_runs(para, rows[r][c], Pt(body_pt), DARK)
    return gf


def _add_footer(slide, sw, sh):
    box = slide.shapes.add_textbox(Inches(0.4), sh - Inches(0.42), sw - Inches(0.8), Inches(0.32))
    p = box.text_frame.paragraphs[0]
    r = p.add_run(); r.text = FOOTER
    r.font.size = Pt(9); r.font.color.rgb = MUTE; p.alignment = PP_ALIGN.RIGHT


def _render_cover(slide, sl, sw, sh):
    bg = slide.shapes.add_shape(1, 0, 0, sw, sh)
    bg.fill.solid(); bg.fill.fore_color.rgb = TEAL; bg.line.fill.background()
    # accent bar
    bar = slide.shapes.add_shape(1, 0, Inches(2.05), sw, Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(0x1F, 0x9D, 0x72); bar.line.fill.background()
    box = slide.shapes.add_textbox(Inches(1.0), Inches(2.4), sw - Inches(2.0), Inches(3.2))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = sl["title"]
    r.font.size = Pt(46); r.font.bold = True; r.font.color.rgb = WHITE
    for kind, val in sl["blocks"]:
        if kind != "p":
            continue
        pp = tf.add_paragraph(); pp.alignment = PP_ALIGN.CENTER; pp.space_before = Pt(12)
        run = pp.add_run(); run.text = _plain(val)
        run.font.size = Pt(18); run.font.color.rgb = CREAM


def build(md_path: str, out_path: str, deck_title: str) -> None:
    base = os.path.dirname(os.path.abspath(md_path))
    with open(md_path, encoding="utf-8") as f:
        slides = parse_slides(f.read())

    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    SW, SH = prs.slide_width, prs.slide_height

    for idx, sl in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        if idx == 0 and sl["level"] == 1:        # cover page
            _render_cover(slide, sl, SW, SH)
            continue

        _add_title(slide, sl["title"] or deck_title, SW)
        tables = [b for b in sl["blocks"] if b[0] == "table"]
        texts = [b for b in sl["blocks"] if b[0] == "p"]
        has_image = sl["image"] and os.path.isfile(os.path.join(base, sl["image"]))
        top = Inches(1.25)

        if tables:
            cur = top
            if texts:
                _add_text(slide, texts, Inches(0.55), cur, SW - Inches(1.1), Inches(0.8))
                cur = cur + Inches(0.7)
            avail = SH - cur - Inches(0.55)
            per = avail / len(tables) - Inches(0.12) if len(tables) > 1 else avail
            for _, rows in tables:
                _add_table(slide, rows, Inches(0.55), cur, SW - Inches(1.1), per)
                cur = cur + per + Inches(0.12)
        elif has_image and texts:
            _add_text(slide, texts, Inches(0.55), top, Inches(6.3), SH - Inches(1.6))
            _place_image(slide, os.path.join(base, sl["image"]), Inches(7.1), Inches(1.4),
                         Inches(5.7), SH - Inches(1.9), SW, center=False)
        elif has_image:
            _place_image(slide, os.path.join(base, sl["image"]), None, Inches(1.4),
                         Inches(10.5), SH - Inches(1.9), SW, center=True)
        else:
            _add_text(slide, texts, Inches(0.55), top, SW - Inches(1.1), SH - Inches(1.6))

        _add_footer(slide, SW, SH)

    prs.save(out_path)
    print(f"✓ Built {out_path} ({len(slides)} slides)")


def _place_image(slide, path, left, top, max_w, max_h, sw, center):
    try:
        pic = slide.shapes.add_picture(path, left or 0, top, width=max_w)
        if pic.height > max_h:
            ratio = pic.width / pic.height
            pic.height = int(max_h); pic.width = int(max_h * ratio)
        if center:
            pic.left = int((sw - pic.width) // 2)
        elif left is not None:
            pic.left = int(left)
    except Exception:
        pass


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("usage: build_pptx.py <input.md> <output.pptx> [deck title]", file=sys.stderr)
        sys.exit(2)
    build(sys.argv[1], sys.argv[2], sys.argv[3] if len(sys.argv) > 3 else "FastClinic")
